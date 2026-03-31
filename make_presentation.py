#!/usr/bin/env python3
"""生成 TeachingMonster TTS 數學公式優化簡報"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# ── 顏色 ──────────────────────────────────────────────
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # 深藍黑
C_ACCENT = RGBColor(0xE5, 0x5D, 0x5D)   # 紅橘
C_LIGHT  = RGBColor(0xF8, 0xF8, 0xF2)   # 米白
C_CODE   = RGBColor(0x2D, 0x2D, 0x44)   # 深灰藍
C_GREEN  = RGBColor(0x50, 0xFA, 0x7B)   # 綠
C_YELLOW = RGBColor(0xFF, 0xB8, 0x6C)   # 橙

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # 空白

# ── 工具函式 ──────────────────────────────────────────
def add_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def textbox(slide, text, left, top, width, height,
            font_size=18, bold=False, color=None, align=PP_ALIGN.LEFT,
            wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    if color:
        run.font.color.rgb = color
    return txBox

def codebox(slide, lines, left, top, width, height, font_size=11):
    """深色背景的程式碼區塊"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.name = "Courier New"
        p.font.color.rgb = RGBColor(0xA9, 0xDC, 0x76)
        p.font.bold = False

def rect(slide, left, top, width, height, fill_color, alpha=None):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def bullet_box(slide, items, left, top, width, height,
               font_size=16, color=None, bullet="●"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet}  {item}"
        p.font.size = Pt(font_size)
        p.space_after = Pt(6)
        if color:
            p.font.color.rgb = color


# ══════════════════════════════════════════════════════
# Slide 1：封面
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_DARK)

rect(s, 0, 0, 13.33, 7.5, RGBColor(0x11, 0x11, 0x22))
rect(s, 0, 0, 0.12, 7.5, C_ACCENT)

textbox(s, "TeachingMonster", 1.2, 1.8, 10, 1.2,
        font_size=52, bold=True, color=C_LIGHT, align=PP_ALIGN.LEFT)
textbox(s, "數學公式朗讀 TTS 優化專案", 1.2, 3.1, 10, 0.9,
        font_size=32, bold=False, color=C_ACCENT, align=PP_ALIGN.LEFT)
textbox(s, "LatexToSpeech 規則引擎的設計與實作", 1.2, 4.1, 9, 0.6,
        font_size=20, color=C_LIGHT, italic=True)

# 右側裝飾
rect(s, 10.2, 1.0, 2.8, 5.5, RGBColor(0xE5, 0x5D, 0x5D))
textbox(s, "TTS", 10.3, 2.8, 2.6, 1.2,
        font_size=56, bold=True, color=C_LIGHT, align=PP_ALIGN.CENTER)
textbox(s, "Qwen3-TTS\n+\nLatexToSpeech", 10.3, 4.0, 2.6, 1.5,
        font_size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

textbox(s, "2026-03-23", 1.2, 6.6, 4, 0.5,
        font_size=14, color=RGBColor(0x88, 0x88, 0x99))


# ══════════════════════════════════════════════════════
# Slide 2：問題背景
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_LIGHT)

rect(s, 0, 0, 13.33, 1.3, C_DARK)
textbox(s, "問題背景", 0.4, 0.3, 5, 0.8,
        font_size=32, bold=True, color=C_LIGHT)
textbox(s, "為什麼需要數學公式朗讀優化？", 0.4, 1.5, 12, 0.6,
        font_size=20, color=RGBColor(0x55, 0x55, 0x66), italic=True)

# 左：問題
rect(s, 0.4, 2.3, 5.8, 4.5, RGBColor(0xFF, 0xF0, 0xF0))
textbox(s, "❌  原始 TTS 的限制", 0.6, 2.5, 5.4, 0.5,
        font_size=18, bold=True, color=C_ACCENT)

problems = [
    "直接餵 LaTeX 符號，變成：\n  x equals left parenthesis...",
    "分數結構成謎：(-b ± √(...))/2a",
    "希臘字母、上下標、積分極限\n  全都破碎不清",
    "學生聽不懂，學習效果大打折扣",
]
bullet_box(s, problems, 0.6, 3.1, 5.4, 3.5, font_size=14, color=C_CODE)

# 右：目標
rect(s, 6.8, 2.3, 5.8, 4.5, RGBColor(0xF0, 0xFF, 0xF0))
textbox(s, "✅  我們的目標", 7.0, 2.5, 5.4, 0.5,
        font_size=18, bold=True, color=RGBColor(0x27, 0xAE, 0x60))

goals = [
    "LaTeX 公式 → 自然口語文字",
    "讓 Qwen3-TTS 朗讀時自然流暢",
    "零額外模型負擔，純規則引擎",
    "可重複使用，無需網路連線",
]
bullet_box(s, goals, 7.0, 3.1, 5.4, 3.5, font_size=14, color=C_CODE)


# ══════════════════════════════════════════════════════
# Slide 3：本來的 TTS 流程
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_DARK)

rect(s, 0, 0, 13.33, 1.3, RGBColor(0x0D, 0x0D, 0x1E))
textbox(s, "原始 TTS 流程", 0.4, 0.3, 8, 0.8,
        font_size=32, bold=True, color=C_LIGHT)
textbox(s, "Qwen3-TTS 的設計假設", 0.4, 1.5, 10, 0.5,
        font_size=18, italic=True, color=RGBColor(0x88, 0x88, 0x99))

# 流程方塊
steps = [
    ("📄 投影片文字", "x = (-b ± √(...))/2a"),
    ("🔧 簡單前處理", "加入句號、簡單替換"),
    ("🤖 Qwen3-TTS", "Ryan 音色 / neutral instruct"),
    ("🔊 音訊輸出", "速度正常但數學破碎"),
]
colors = [RGBColor(0x34,0x59,0x8B), RGBColor(0x3D,0x6D,0x4A),
          RGBColor(0x8B,0x3D,0x5D), RGBColor(0x5D,0x3D,0x8B)]
for i, (title, detail) in enumerate(steps):
    x = 0.5 + i * 3.15
    rect(s, x, 2.4, 2.8, 2.8, colors[i])
    textbox(s, title, x+0.1, 2.7, 2.6, 0.5,
            font_size=16, bold=True, color=C_LIGHT, align=PP_ALIGN.CENTER)
    textbox(s, detail, x+0.1, 3.4, 2.6, 1.4,
            font_size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)
    if i < 3:
        textbox(s, "→", x+2.75, 3.4, 0.4, 0.5,
                font_size=28, color=C_ACCENT, align=PP_ALIGN.CENTER)

# 下方說明
rect(s, 0.5, 5.5, 12.3, 0.05, C_ACCENT)
textbox(s, "Qwen3-TTS 專為自然語言語音合成設計，吃乾淨的文字，輸出流暢語音。\n"
           "但數學公式是另一種語言——它需要被「翻譯」成人類朗讀的樣子。",
        0.5, 5.7, 12, 1.2,
        font_size=15, color=RGBColor(0xBB, 0xBB, 0xCC), italic=True)


# ══════════════════════════════════════════════════════
# Slide 4：研究方案
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_LIGHT)

rect(s, 0, 0, 13.33, 1.3, C_DARK)
textbox(s, "研究與方案", 0.4, 0.3, 8, 0.8,
        font_size=32, bold=True, color=C_LIGHT)
textbox(s, "查詢到的三種可行方法", 0.4, 1.5, 10, 0.5,
        font_size=18, italic=True, color=RGBColor(0x88, 0x88, 0x99))

cards = [
    {
        "title": "MathReader (2025 新論文)",
        "desc": "Nougat OCR + T5-small 翻譯 + VITS",
        "pros": "最新研究，WER 大幅降低",
        "cons": "需訓練模型，GPU 負擔大",
        "color": RGBColor(0x34, 0x59, 0x8B),
        "tag": "論文方案",
    },
    {
        "title": "Speech-Rule-Engine",
        "desc": "JavaScript 規則引擎",
        "pros": "速度快，有標準規範",
        "cons": "需跑 JS 環境，整合困難",
        "color": RGBColor(0x27, 0xAE, 0x60),
        "tag": "國外開源",
    },
    {
        "title": "LatexToSpeech 規則引擎",
        "desc": "Python 正則表達式 + 對照表",
        "pros": "完全可控、輕量、無網路依賴",
        "cons": "複雜巢狀公式需要持續擴充",
        "color": C_ACCENT,
        "tag": "我們選用的方案 ✅",
    },
]

for i, c in enumerate(cards):
    x = 0.4 + i * 4.3
    rect(s, x, 2.2, 4.0, 4.6, c["color"])
    # tag
    tag_color = RGBColor(0xFF, 0xFF, 0xFF) if c["color"] != C_ACCENT else RGBColor(0xFF, 0xFF, 0xCC)
    textbox(s, c["tag"], x+0.1, 2.35, 3.8, 0.35,
            font_size=10, bold=True, color=tag_color)
    # title
    textbox(s, c["title"], x+0.15, 2.8, 3.7, 0.6,
            font_size=16, bold=True, color=RGBColor(0xFF, 0xFF, 0xFF))
    # desc
    textbox(s, c["desc"], x+0.15, 3.45, 3.7, 0.5,
            font_size=12, color=RGBColor(0xDD, 0xDD, 0xDD))
    # pros
    textbox(s, f"✅ {c['pros']}", x+0.15, 4.1, 3.7, 0.6,
            font_size=12, color=RGBColor(0xA0, 0xFF, 0xA0))
    # cons
    textbox(s, f"⚠️  {c['cons']}", x+0.15, 4.75, 3.7, 0.7,
            font_size=12, color=RGBColor(0xFF, 0xCC, 0xAA))
    # arrow
    if i < 2:
        textbox(s, "→", x+3.95, 4.2, 0.4, 0.5,
                font_size=28, color=RGBColor(0x88,0x88,0x88), align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════
# Slide 5：LatexToSpeech 架構
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_DARK)

rect(s, 0, 0, 13.33, 1.3, RGBColor(0x0D, 0x0D, 0x1E))
textbox(s, "LatexToSpeech 規則引擎", 0.4, 0.3, 10, 0.8,
        font_size=32, bold=True, color=C_LIGHT)
textbox(s, "插入在 Qwen3-TTS 之前，純 Python 輕量前處理", 0.4, 1.5, 10, 0.5,
        font_size=18, italic=True, color=RGBColor(0x88,0x88,0x99))

# 流程
pipeline = [
    ("1. 輸入", "含 LaTeX 的文字", RGBColor(0x34,0x59,0x8B)),
    ("2. 分數", "\\frac{a}{b}", RGBColor(0x8B,0x3D,0x5D)),
    ("3. 極限/總和", "\\lim \\sum \\int", RGBColor(0x5D,0x8B,0x3D)),
    ("4. 上標/下標", "x^2, x_i", RGBColor(0x8B,0x5D,0x3D)),
    ("5. 根號/函數", "\\sqrt, \\sin", RGBColor(0x3D,0x5D,0x8B)),
    ("6. 符號/字母", "\\alpha, \\leq", RGBColor(0x8B,0x8B,0x3D)),
    ("7. 清理輸出", "乾淨口語文字", RGBColor(0x59,0x8B,0x3D)),
]

for i, (step, example, col) in enumerate(pipeline):
    x = 0.35 + i * 1.83
    rect(s, x, 2.3, 1.65, 1.7, col)
    textbox(s, step, x+0.05, 2.4, 1.55, 0.4,
            font_size=13, bold=True, color=C_LIGHT, align=PP_ALIGN.CENTER)
    textbox(s, example, x+0.05, 2.9, 1.55, 0.9,
            font_size=10, color=RGBColor(0xCC,0xCC,0xDD), align=PP_ALIGN.CENTER)
    if i < 6:
        textbox(s, "→", x+1.58, 2.9, 0.25, 0.4,
                font_size=18, color=C_ACCENT, align=PP_ALIGN.CENTER)

# 支援項目
textbox(s, "支援的轉換項目", 0.4, 4.3, 5, 0.4,
        font_size=15, bold=True, color=C_ACCENT)
items = [
    "分數 \\frac{a}{b}  →  the fraction with numerator a, and denominator b",
    "上標 x^{2}         →  x squared / x to the power of n",
    "下標 x_{ij}        →  x sub i j",
    "根號 \\sqrt{x}     →  the square root of x",
    "希臘字母 \\alpha    →  alpha",
    "極限 \\lim_{x→0}   →  limit as x approaches 0",
    "總和 \\sum_{i=1}^{n} →  sum from i equals 1 to n",
    "積分 \\int_0^1      →  integral from 0 to 1",
    "矩陣 \\begin{pmatrix} a & b \\\\ c & d \\end{pmatrix}",
]
bullet_box(s, items, 0.4, 4.75, 6.3, 2.5, font_size=11, color=RGBColor(0xAA,0xCC,0xAA))

# 右側：核心對照表
rect(s, 7.0, 4.3, 5.8, 3.0, RGBColor(0x22,0x22,0x38))
textbox(s, "符號對照表（精選）", 7.15, 4.4, 5.5, 0.4,
        font_size=14, bold=True, color=C_ACCENT)
table_lines = [
    "\\\\frac{a}{b}  →  a over b",
    "\\\\sqrt{x}     →  square root of x",
    "\\\\alpha       →  alpha",
    "\\\\leq         →  less than or equal to",
    "\\\\times      →  times",
    "\\\\rightarrow →  goes to",
    "\\\\sin         →  sine",
    "\\\\int_0^1    →  integral from 0 to 1",
]
for j, line in enumerate(table_lines):
    textbox(s, line, 7.15, 4.9 + j * 0.27, 5.5, 0.28,
            font_size=11, color=RGBColor(0xA9,0xDC,0x76), italic=False)


# ══════════════════════════════════════════════════════
# Slide 6：轉換範例
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_LIGHT)

rect(s, 0, 0, 13.33, 1.3, C_DARK)
textbox(s, "實際轉換範例", 0.4, 0.3, 8, 0.8,
        font_size=32, bold=True, color=C_LIGHT)
textbox(s, "從 LaTeX 到流暢朗讀文字", 0.4, 1.5, 10, 0.5,
        font_size=18, italic=True, color=RGBColor(0x88,0x88,0x99))

examples = [
    {
        "latex": r"x = \frac{-b + \sqrt{b^2-4ac}}{2a}",
        "before": "x equals left parenthesis negative b plus or minus root...",
        "after":  "x equals the fraction with numerator minus b plus the square root of b squared minus 4 a c, and denominator 2 a",
        "tag": "一元二次公式",
    },
    {
        "latex": r"e^{i\pi} + 1 = 0",
        "before": "e left superscript i pi right superscript plus 1 equals 0",
        "after":  "e to the power of i pi plus 1 equals 0",
        "tag": "歐拉公式",
    },
    {
        "latex": r"\int_0^1 x^2 dx = \frac{1}{3}",
        "before": "integral sub 0 to the power of one x squared dx...",
        "after":  "integral from 0 to 1 x squared dx equals the fraction with numerator 1, and denominator 3",
        "tag": "積分計算",
    },
    {
        "latex": r"\sum_{n=1}^{\infty} \frac{1}{n^2}",
        "before": "sum sub n equals 1 to the power of infinity...",
        "after":  "sum from n equals 1 to infinity the fraction with numerator 1, and denominator n squared",
        "tag": "級數",
    },
]

for i, ex in enumerate(examples):
    row = i // 2
    col = i % 2
    x = 0.4 + col * 6.5
    y = 2.1 + row * 2.55

    rect(s, x, y, 6.2, 2.35, RGBColor(0xF5,0xF0,0xFF))
    # tag
    textbox(s, ex["tag"], x+0.1, y+0.1, 1.5, 0.3,
            font_size=10, bold=True, color=RGBColor(0x88,0,0))
    # latex
    textbox(s, f"LaTeX：{ex['latex']}", x+0.1, y+0.45, 6.0, 0.45,
            font_size=10, color=RGBColor(0x33,0x33,0x66))
    # before
    textbox(s, f"❌ {ex['before']}", x+0.1, y+0.95, 6.0, 0.55,
            font_size=10, color=RGBColor(0x99,0,0))
    # after
    textbox(s, f"✅ {ex['after']}", x+0.1, y+1.55, 6.0, 0.7,
            font_size=10, color=RGBColor(0x00,0x66,0x00))


# ══════════════════════════════════════════════════════
# Slide 7：整合方式
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_DARK)

rect(s, 0, 0, 13.33, 1.3, RGBColor(0x0D,0x0D,0x1E))
textbox(s, "與 TeachingMonster 整合", 0.4, 0.3, 10, 0.8,
        font_size=32, bold=True, color=C_LIGHT)
textbox(s, "一行程式碼升級現有流程", 0.4, 1.5, 10, 0.5,
        font_size=18, italic=True, color=RGBColor(0x88,0x88,0x99))

# 整合示意圖
rect(s, 0.4, 2.2, 2.8, 1.0, RGBColor(0x34,0x59,0x8B))
textbox(s, "原始投影片文字", 0.5, 2.35, 2.6, 0.7,
        font_size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)
textbox(s, "↓ + LatexToSpeech", 0.5, 3.3, 2.6, 0.4,
        font_size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)

rect(s, 0.4, 3.7, 2.8, 1.0, RGBColor(0x50,0xFA,0x7B))
textbox(s, "口語化文字", 0.5, 3.85, 2.6, 0.7,
        font_size=13, color=C_DARK, align=PP_ALIGN.CENTER, bold=True)

textbox(s, "↓ Qwen3-TTS", 0.5, 4.8, 2.6, 0.4,
        font_size=12, color=C_YELLOW, align=PP_ALIGN.CENTER)

rect(s, 0.4, 5.2, 2.8, 1.0, RGBColor(0x8B,0x3D,0x5D))
textbox(s, "高品質語音輸出", 0.5, 5.35, 2.6, 0.7,
        font_size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

# 程式碼範例
rect(s, 3.8, 2.2, 8.8, 5.0, RGBColor(0x1A,0x1A,0x2E))
textbox(s, "使用範例", 4.0, 2.35, 3, 0.4,
        font_size=14, bold=True, color=C_ACCENT)

code = [
    "from src.tts.latex_to_speech import LatexToSpeech",
    "from src.tts.tts import TTSModule",
    "",
    "# 初始化",
    "latex2speech = LatexToSpeech()",
    "tts = TTSModule(...)",
    "tts.load()",
    "",
    "# 投影片內容",
    "slides = [",
    '    r\"x = \\\\frac{{-b + \\\\sqrt{{b^2-4ac}}}}{{2a}}\",',
    '    r\"\\\\int_0^1 x^2 dx = \\\\frac{{1}}{{3}}\",',
    "]",
    "",
    "# 轉換 + 朗讀",
    "spoken = [latex2speech.convert_inline(s) for s in slides]",
    "tts.run(spoken)",
]
codebox(s, code, 4.0, 2.85, 8.3, 4.2, font_size=11)


# ══════════════════════════════════════════════════════
# Slide 8：修改記錄
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_LIGHT)

rect(s, 0, 0, 13.33, 1.3, C_DARK)
textbox(s, "修改記錄", 0.4, 0.3, 8, 0.8,
        font_size=32, bold=True, color=C_LIGHT)
textbox(s, "今天在 TeachingMonster 做的所有變更", 0.4, 1.5, 12, 0.5,
        font_size=18, italic=True, color=RGBColor(0x88,0x88,0x99))

changes = [
    ("🆕 新增檔案", "src/tts/latex_to_speech.py",
     "23 KB 規則引擎，涵蓋分數/次方/根號/極限/總和/積分/矩陣等"),
    ("🐛 修復矩陣匹配順序",
     "_convert_matrix 在 _convert_envs 之後執行",
     "導致 pmatrix begin/end 被移除後才匹配 → 改順序後正常"),
    ("🐛 修復積分邊界",
     "\\int_0^1 被錯誤解析為 sub 0 上標 power of 1",
     "加入 pattern6/7 專門處理無括號積分邊界"),
    ("📋 新增測試腳本", "test_math.py",
     "用 conda monster 環境驗證 6 段數學公式朗讀正確生成"),
    ("🔧 執行環境確認", "使用 conda run -n monster",
     "Qwen3-TTS + faster-whisper 正確在 monster 環境執行"),
    ("📄 新增簡報", "test_math_presentation.pptx",
     "本份簡報，記錄今日所有工作成果"),
]

for i, (tag, title, detail) in enumerate(changes):
    y = 2.1 + i * 0.88
    rect(s, 0.4, y, 12.5, 0.78, RGBColor(0xF0,0xF4,0xFF))
    textbox(s, tag, 0.55, y+0.05, 1.6, 0.32,
            font_size=11, bold=True, color=RGBColor(0x33,0x33,0x99))
    textbox(s, title, 2.2, y+0.05, 4.5, 0.32,
            font_size=13, bold=True, color=RGBColor(0x22,0x22,0x44))
    textbox(s, detail, 2.2, y+0.4, 10.5, 0.35,
            font_size=11, color=RGBColor(0x55,0x55,0x77))


# ══════════════════════════════════════════════════════
# Slide 9：限制與未來方向
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_DARK)

rect(s, 0, 0, 13.33, 1.3, RGBColor(0x0D,0x0D,0x1E))
textbox(s, "限制與未來方向", 0.4, 0.3, 10, 0.8,
        font_size=32, bold=True, color=C_LIGHT)

# 左：限制
rect(s, 0.4, 1.8, 5.8, 5.0, RGBColor(0x22,0x22,0x38))
textbox(s, "⚠️  目前限制", 0.6, 2.0, 5.4, 0.45,
        font_size=18, bold=True, color=C_ACCENT)
limits = [
    "複雜巢狀分數可能需要多次迭代轉換",
    "自定義 LaTeX 巨集（\\def, \\newcommand）不支援",
    "矩陣複雜度僅支援到 4x4",
    "無 Anselm 書籍那種特殊數學朗讀慣例",
    "沒有針對 Taylor/Maclaurin 展開的特殊處理",
]
bullet_box(s, limits, 0.6, 2.55, 5.4, 4.0, font_size=14, color=RGBColor(0xBB,0xBB,0xDD))

# 右：方向
rect(s, 6.8, 1.8, 5.8, 5.0, RGBColor(0x22,0x22,0x38))
textbox(s, "🔮  未來方向", 7.0, 2.0, 5.4, 0.45,
        font_size=18, bold=True, color=C_GREEN)
directions = [
    "用 MathReader 的 T5-small 微調模型替換規則引擎",
    "加入 MathBridge 資料集訓練",
    "支援中文數學朗讀（\\frac →分之）",
    "整合 Edge TTS 作為備用引擎",
    "自動偵測 \\begin{equation} 環境",
    "加進 TeachingMonster 的 slides 流程中",
]
bullet_box(s, directions, 7.0, 2.55, 5.4, 4.0, font_size=14, color=RGBColor(0xBB,0xBB,0xDD))

# ══════════════════════════════════════════════════════
# Slide 10：總結
# ══════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_bg(s, C_DARK)

rect(s, 0, 0, 13.33, 7.5, RGBColor(0x11,0x11,0x22))
rect(s, 0, 0, 0.12, 7.5, C_ACCENT)

textbox(s, "總結", 1.2, 1.0, 10, 0.8,
        font_size=40, bold=True, color=C_LIGHT)
rect(s, 1.2, 1.85, 10, 0.05, C_ACCENT)

summary = [
    "✅  建構了 LatexToSpeech 規則引擎，純 Python、輕量、無依賴",
    "✅  成功轉換複雜 LaTeX 數學公式為自然口語文字",
    "✅  與 Qwen3-TTS 整合，朗讀數學公式流暢自然",
    "✅  支援分數、根號、上標、下標、極限、總和、積分、矩陣等",
    "✅  測試驗證成功，mp3 正常生成",
]
bullet_box(s, summary, 1.2, 2.1, 11, 3.5, font_size=18, color=RGBColor(0xAA,0xDD,0xAA))

rect(s, 1.2, 5.3, 10, 0.05, C_ACCENT)
textbox(s, "這只是起點。\n"
           "下一步：T5 微調模型，把規則引擎升級成 ML 模型，\n"
           "讓每一個數學式都朗讀精準。",
        1.2, 5.5, 10, 1.5,
        font_size=16, italic=True, color=RGBColor(0x88,0x88,0x99))

textbox(s, "檔案位置：TeachingMonster-released/src/tts/latex_to_speech.py",
        1.2, 7.0, 10, 0.4, font_size=13, color=RGBColor(0x55,0x55,0x77))


# ── 儲存 ──────────────────────────────────────────────
output_path = "TeachingMonster-released/tts_math_presentation.pptx"
prs.save(output_path)
print(f"簡報已儲存：{output_path}")
